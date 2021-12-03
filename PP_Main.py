import os
import sys
from PyQt5.QtWidgets import *
from PyQt5.QtGui import *
from PyQt5.QtCore import QDate, QLocale, Qt, QModelIndex, QSize
from PyQt5 import uic
from openpyxl import load_workbook
from openpyxl.styles import Border, Side, PatternFill, Font
from openpyxl.utils.cell import coordinate_from_string, column_index_from_string
import shutil
import json
from subprocess import Popen, PIPE

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor


os.environ['QT_MAC_WANTS_LAYER'] = '1' #필수
myClasses =[]
myStudents = []
myFiles = []
constDay = ['월', '화', '수', '목', '금', '토', '일']
with open(os.getcwd() + '/savedata.json', 'r',encoding='UTF-8') as f:
    json_data = json.load(f)
config = json_data

with open(os.getcwd() + '/path.json', 'r',encoding='UTF-8') as f:
    path_data = json.load(f)
pathConfig = path_data
    
studentListModel = QStandardItemModel()
classListModel = QStandardItemModel()
fileListModel = QStandardItemModel()
kakaotalkScript = '''
on run {x,y}
tell application "KakaoTalk"
	reopen -- unminimizes the first minimized window or makes a new default window
	activate (first window whose name is "카카오톡") -- makes the app frontmost
end tell

tell application "System Events" to key code 18 using command down -- activate main window
delay 2
tell application "System Events" to key code 3 using command down
delay 2
set the clipboard to x
tell application "System Events" to key code 9 using command down
delay 2
tell application "System Events" to key code 125
tell application "System Events" to key code 36

set the clipboard to y
tell application "System Events" to key code 9 using command down
tell application "System Events" to key code 36
end run

on splitText(theText, theDelimiter)
        set AppleScript's text item delimiters to theDelimiter
        set theTextItems to every text item of theText
        set AppleScript's text item delimiters to ""
        return theTextItems
end splitText
'''
form_class = uic.loadUiType(os.getcwd() + "/PersonalProject.ui")[0]

class MyWindow(QMainWindow, form_class):

    def __init__(self):
        global myClasses
        global config
        super().__init__()
        
        self.setWindowTitle("메이킷코드랩")
        self.setWindowIcon(QIcon(os.getcwd() + "/PP_icon.jpg"))


        self.setupUi(self)

        myClasses = list(config.keys())
        for x in myClasses:
            classListModel.appendRow(QStandardItem(x))

        self.feedbackPathField.setText(pathConfig.get('feedbackPath', ''))
        self.jindoPathField.setText(pathConfig.get('jindoPath', ''))
        self.dailyReportPathField.setText(pathConfig.get('dailyReportPath', ''))
        self.NASPathField.setText(pathConfig.get('NASPath', ''))
        

        self.classList.setModel(classListModel)
        self.studentList.setModel(studentListModel)
        self.fileList.setModel(fileListModel)
        
        self.classList.clicked.connect(self.classSelectHandler)
        self.studentList.clicked.connect(self.studentSelectHandler)

        self.feedbackField.textChanged.connect(self.feedbackFieldHandler)

        self.classCommentField.textChanged.connect(self.classCommentFieldHandler)
        self.classSpecialField.textChanged.connect(self.classSpecialFieldHandler)
        self.classHomeworkField.textChanged.connect(self.classHomeworkFieldHandler)

        self.noHomeworkCheck.stateChanged.connect(self.noHomeworkCheckHandler)

        self.addFileButton.clicked.connect(self.addFileButtonHandler)
        self.removeFileButton.clicked.connect(self.removeFileButtonHandler)

        self.isHome.stateChanged.connect(self.isHomeHandler)
        self.noKakaoCheck.clicked.connect(self.noKakaoCheckHandler)
        self.kakaoNameField.textChanged.connect(self.kakaoNameFieldHandler)

        self.setFeedbackPathButton.clicked.connect(self.setFeedbackPathButtonHandler)
        self.setJindoPathButton.clicked.connect(self.setJindoPathButtonHandler)
        self.setDailyReportPathButton.clicked.connect(self.setDailyReportPathButtonHandler)
        self.setNASPathButton.clicked.connect(self.setNASPathButtonHandler)

        self.saveFeedback.clicked.connect(self.saveFeedbackHandler)
        self.sendKakaoButton.clicked.connect(self.allStudentSendKakao)
        self.saveExcel.clicked.connect(self.autoxl)
        self.savePPT.clicked.connect(self.autoppt)
        self.saveReadmeButton.clicked.connect(self.saveReadme)

        self.addClassButton.clicked.connect(self.addClassHandler)
        self.removeClassButton.clicked.connect(self.removeClassHandler)
        self.classColField.textChanged.connect(self.classColFieldHandler)

        self.addStudentButton.clicked.connect(self.addStudentHandler)
        self.removeStudentButton.clicked.connect(self.removeStudentHandler)


        self.inquiry() #statusBar에 시간 출력하기

    def addFileButtonHandler(self):
        global myFiles
        global config

        fileName, _ = QFileDialog.getOpenFileName(self, '파일 선택', '', 'All Files (*)')
        if fileName:
            myFiles.append(fileName)
            fileListModel.appendRow(QStandardItem(fileName))
            config[self.getCurrentClass()]['files']=myFiles
            self.fileList.setModel(fileListModel)

    def removeFileButtonHandler(self):
        global myFiles
        global config

        index = self.fileList.currentIndex()
        if index.isValid():
            fileListModel.removeRow(index.row())
            myFiles.remove(myFiles[index.row()])
            config[self.getCurrentClass()]['files']=myFiles
            self.fileList.setModel(fileListModel)

    def noHomeworkCheckHandler(self):
        global config
        config[self.getCurrentClass()]['noHomework'] = self.noHomeworkCheck.isChecked()

    def saveReadme(self):
        homePeople = []
        for student in config[self.getCurrentClass()]['students']:
                if config[self.getCurrentClass()]['students'][student]['isHome']:
                    homePeople.append(student + " 결석")
        readme = "진도 : " + config[self.getCurrentClass()]['classComment'] + "\n" \
                + "결석 : " + ", ".join(homePeople) + "\n" 
        if not config[self.getCurrentClass()].get('noHomework', False):
                readme += "숙제 : " + config[self.getCurrentClass()]['classHomework'] + "\n"

        reportDir = os.getcwd() + "/report/" + config[self.getCurrentClass()]['folderName']

        if not os.path.exists(reportDir):
            os.makedirs(reportDir)

        #initialize
        for file in os.scandir(reportDir):
            os.remove(file.path)
        
        with open(reportDir+'/readme.txt', 'w+', encoding='UTF-8') as f:
            f.write(readme)

        #save file
        for file in myFiles:
            shutil.copy(file, reportDir)
        
    def addClassHandler(self):
        uclassName, ok = QInputDialog.getText(self, '반 추가', '반 이름 입력:')
        if ok:
            config[uclassName] = {'students':{}, \
                                    'classSpecial':'', \
                                    'classComment':'', \
                                    "folderName": "___",
                                    "classComment": "",
                                    "classSpecial": "",
                                    "excelCol": 4
            }
            myClasses.append(uclassName)
            classListModel.appendRow(QStandardItem(uclassName))

            classFolderName, fok = QInputDialog.getText(self, '폴더 이름 설정', '수업 자료 폴더 이름 입력:')
            if fok:
                config[uclassName]['folderName'] = classFolderName
    
    def removeClassHandler(self):
        config.pop(self.getCurrentClass())
        myClasses.remove(self.getCurrentClass())
        classListModel.removeRow(self.classList.currentIndex().row())          

    def classColFieldHandler(self):
        global config
        config[self.getCurrentClass()]['excelCol'] = self.classColField.text()

    def removeStudentHandler(self):
        config[self.getCurrentClass()]['students'].pop(self.getCurrentStudent())
        myStudents.remove(self.getCurrentStudent())
        studentListModel.removeRow(self.studentList.currentIndex().row())
    
    def addStudentHandler(self):
        studentName, ok = QInputDialog.getText(self, '학생 추가', '학생 이름 입력:')
        if ok:
            config[self.getCurrentClass()]['students'][studentName] = {
				"isHome": False,
				"feedback": "",
				"kakaoName": "박우현",
				"noKakao": False
                }
            studentListModel.appendRow(QStandardItem(studentName))
            myStudents.append(studentName)

    def noKakaoCheckHandler(self):
        global config
        if self.noKakaoCheck.isChecked():
            config[self.getCurrentClass()]['students'][self.getCurrentStudent()]['noKakao'] = True
        else:
            config[self.getCurrentClass()]['students'][self.getCurrentStudent()]['noKakao'] = False

    def kakaoNameFieldHandler(self):
        global config
        config[self.getCurrentClass()]['students'][self.getCurrentStudent()]['kakaoName'] = self.kakaoNameField.text()
        
    def setFeedbackPathButtonHandler(self):
        global pathConfig
        feedbackPath, ok = QFileDialog.getOpenFileName(self, '파일 선택', '', 'All Files (*)')
        if ok:
            pathConfig['feedbackPath'] = feedbackPath
            self.feedbackPathField.setText(feedbackPath)

    def setJindoPathButtonHandler(self):
        global pathConfig
        jindoPath, ok = QFileDialog.getOpenFileName(self, '파일 선택', '', 'All Files (*)')
        if ok:
            pathConfig['jindoPath'] = jindoPath
            self.jindoPathField.setText(jindoPath)
    
    def setDailyReportPathButtonHandler(self):
        global pathConfig
        dailyReportPath, ok = QFileDialog.getOpenFileName(self, '파일 선택', '', 'All Files (*)')
        if ok:
            pathConfig['dailyReportPath'] = dailyReportPath
            self.dailyReportPathField.setText(dailyReportPath)
    
    def setNASPathButtonHandler(self):
        global pathConfig
        NASPath= QFileDialog.getExistingDirectory(self, '폴더 선택', '', QFileDialog.ShowDirsOnly)
        ok = True
        if ok:
            pathConfig['NASPath'] = NASPath
            self.NASPathField.setText(NASPath)

    def isHomeHandler(self):
        if self.isHome.isChecked():
            config[self.getCurrentClass()]['students'][self.getCurrentStudent()]['isHome'] = True
        else:
            config[self.getCurrentClass()]['students'][self.getCurrentStudent()]['isHome'] = False

    def autoppt(self):
        QLocale.setDefault(QLocale(QLocale.Korean, QLocale.SouthKorea))

        
        cur_date = QDate.currentDate()
        str_date = cur_date.toString("MM/dd")
        str_date += '('+constDay[cur_date.dayOfWeek()-1]+')'
        #print(str_date)
        prs = Presentation(pathConfig.get('dailyReportPath', 'dailyReport.pptx'))
        '''
        for i in range(0, 11):
            print("--------[%d] ------ "%(i))
            slide = prs.slides.add_slide(prs.slide_layouts[i])
            for shape in slide.placeholders:
                print('%d %s' % (shape.placeholder_format.idx, shape.name))
        '''

        #16 = Date, 15, 14, 13,..
        slide_layout = prs.slide_layouts[7]
        slide = prs.slides.add_slide(slide_layout)
        tf = slide.placeholders[16].text_frame
        tf.text =str_date

        #출결
        idx=-1
        for i in range(0, len(myClasses)):
            myClassName = myClasses[i]
            if not constDay[cur_date.dayOfWeek()-1] in myClassName:
                continue
            idx+=1
            homePeople = []
            for student in config[myClassName]['students']:
                if config[myClassName]['students'][student]['isHome']:
                    homePeople.append(student + " 결석")
            if len(homePeople) == 0:
                homePeople.append('전원출석')
            homePeopleStr = ','.join(homePeople)

            #comment, special
            commentStr = config[myClassName]['classComment']
            specialStr = config[myClassName]['classSpecial']
            if specialStr == '':
                specialStr = '없습니다.'

            tf = slide.placeholders[15-idx].text_frame
            tf.text = '반이름 : ' + \
                myClassName + '\n' + \
                    '출결 : ' + homePeopleStr + '\n' + \
                    '진도 : ' + commentStr + '\n' + \
                    '특이사항 : '+specialStr

        self.move_slide(prs, prs.slides.index(slide), 1)

        prs.save(pathConfig.get('dailyReportPath', 'dailyReport.pptx'))

    def move_slide(self, presentation, old_index, new_index):
        xml_slides = presentation.slides._sldIdLst  # pylint: disable=W0212
        slides = list(xml_slides)
        xml_slides.remove(slides[old_index])
        xml_slides.insert(new_index, slides[old_index])

    def classSpecialFieldHandler(self):
        config[self.getCurrentClass()]['classSpecial'] = self.classSpecialField.toPlainText()

    def classCommentFieldHandler(self):
        config[self.getCurrentClass()]['classComment'] = self.classCommentField.toPlainText()

    def classHomeworkFieldHandler(self):
        config[self.getCurrentClass()]['classHomework'] = self.classHomeworkField.toPlainText()

    def allStudentSendKakao(self):
        for x in myStudents:
            if config[self.getCurrentClass()]['students'][x]['isHome'] or config[self.getCurrentClass()]['students'][x]['noKakao']:
                continue
            self.sendKakaotalk(config[self.getCurrentClass()]['students'][x]['kakaoName'],\
            config[self.getCurrentClass()]['students'][x]['feedback'])
    
    def sendKakaotalk(self, x, y):
        p = Popen(['osascript', '-'] + [x,y], stdin=PIPE, stdout=PIPE, stderr=PIPE)
        stdout, stderr = p.communicate(kakaotalkScript.encode('utf-8'))

    def saveFeedbackHandler(self):
        with open(os.getcwd() + '/savedata.json', 'w',encoding='UTF-8') as f:
            json.dump(config, f, ensure_ascii=False)
        
        with open(os.getcwd() + '/path.json', 'w',encoding='UTF-8') as f:
            json.dump(pathConfig, f, ensure_ascii=False)

    def feedbackFieldHandler(self):
        config[self.getCurrentClass()]['students'][self.getCurrentStudent()]['feedback'] = self.feedbackField.toPlainText()

    def getCurrentClass(self):
        return myClasses[self.classList.currentIndex().row()]
    
    def getCurrentStudent(self):
        return myStudents[self.studentList.currentIndex().row()]

    def studentSelectHandler(self):
        self.feedbackField.setPlainText(config[self.getCurrentClass()]['students'][self.getCurrentStudent()]['feedback'])

        #load kakao
        self.kakaoNameField.setText(config[self.getCurrentClass()]['students'][self.getCurrentStudent()]['kakaoName'])

        #load isHome
        if config[self.getCurrentClass()]['students'][self.getCurrentStudent()]['isHome'] == True:
            self.isHome.setChecked(True)
        else:
            self.isHome.setChecked(False)

        #load noKakao
        if config[self.getCurrentClass()]['students'][self.getCurrentStudent()]['noKakao'] == True:
            self.noKakaoCheck.setChecked(True)
        else:
            self.noKakaoCheck.setChecked(False)

    def classSelectHandler(self):
        global myStudents
        global myFiles
        myStudents =list(config[myClasses[self.classList.currentIndex().row()]]['students'].keys())
        studentListModel.clear()
        for x in myStudents:
            studentListModel.appendRow(QStandardItem(x))
        
        #load files
        myFiles = config[self.getCurrentClass()].get('files', [])
        fileListModel.clear()
        for x in myFiles:
            fileListModel.appendRow(QStandardItem(x))
        

        #load comment, special
        self.classSpecialField.setPlainText(config[self.getCurrentClass()]['classSpecial'])
        self.classCommentField.setPlainText(config[self.getCurrentClass()]['classComment'])
        self.classHomeworkField.setPlainText(config[self.getCurrentClass()].get('classHomework', ''))
        self.noHomeworkCheck.setChecked(config[self.getCurrentClass()].get('noHomework', False))
        self.classColField.setText(str(config[self.getCurrentClass()].get('excelCol', '')))

    def autoxl(self):
        global myClasses, myStudents
        try:
            wb = load_workbook(filename=pathConfig.get('feedbackPath', 'feedback.xlsx'))
            tmpClassName = self.getCurrentClass().replace('/',',')
            tmpClassName = tmpClassName.replace(':','시')
            ws = wb[tmpClassName]
            row = 2
            col = 2

            while (ws.cell(row=row, column=col).value != None):
                row += 2
            
            QLocale.setDefault(QLocale(QLocale.Korean, QLocale.SouthKorea))
            cur_date = QDate.currentDate()
            str_date = cur_date.toString(Qt.DefaultLocaleLongDate)

            titleColor = PatternFill(start_color='EDEDED', end_color='EDEDED', fill_type='solid')
            valueColor = PatternFill(start_color='FFFFFF', end_color='FFFFFF', fill_type='solid')
            defaultFont = Font(name='맑은 고딕', size=9, bold=False, italic=False, color='000000')


            ws.cell(row=row, column=col).value = str_date

            #fill
            ws.cell(row=row,column=col).fill = titleColor
            ws.cell(row=row+1,column=col).fill = valueColor

            #font
            ws.cell(row=row,column=col).font = defaultFont
            ws.cell(row=row+1,column=col).font = defaultFont

            #border
            ws.cell(row=row, column=col).border = Border(left=Side(border_style='thin', color='000000'),\
                                                                right=Side(border_style='thin', color='000000'),\
                                                                top=Side(border_style='thin', color='000000'),\
                                                                bottom=Side(border_style='thin', color='000000'))
            ws.cell(row=row+1, column=col).border = Border(left=Side(border_style='thin', color='000000'),\
                                                                right=Side(border_style='thin', color='000000'),\
                                                                top=Side(border_style='thin', color='000000'),\
                                                                bottom=Side(border_style='thin', color='000000'))
            
            for i in range(len(config[self.getCurrentClass()]['students'].keys())):
                #border
                ws.cell(row=row, column=col+i+1).border = Border(left=Side(border_style='thin', color='000000'),\
                                                                right=Side(border_style='thin', color='000000'),\
                                                                top=Side(border_style='thin', color='000000'),\
                                                                bottom=Side(border_style='thin', color='000000'))
                ws.cell(row=row+1, column=col+i+1).border = Border(left=Side(border_style='thin', color='000000'),\
                                                                right=Side(border_style='thin', color='000000'),\
                                                                top=Side(border_style='thin', color='000000'),\
                                                                bottom=Side(border_style='thin', color='000000'))
            
                #font
                ws.cell(row=row,column=col+i+1).font = defaultFont
                ws.cell(row=row+1,column=col+i+1).font = defaultFont 

                #fill
                ws.cell(row=row,column=col+i+1).fill = titleColor
                ws.cell(row=row+1,column=col+i+1).fill = valueColor

                #value
                ws.cell(row=row,column=col+i+1).value = myStudents[i]
                if(config[self.getCurrentClass()]['students'][myStudents[i]]['isHome'] == True):
                    ws.cell(row=row+1,column=col+i+1).value = '결석'
                else:
                    ws.cell(row=row+1,column=col+i+1).value = config[self.getCurrentClass()]['students'][myStudents[i]]['feedback']

            wb.save(pathConfig.get('feedbackPath', 'feedback.xlsx'))
        except:
            print('error')

        # jindo Excel
        wb = load_workbook(filename=pathConfig.get('jindoPath',''))
        defaultFont = Font(name='맑은 고딕', size=11, bold=False, italic=False, color='000000')
        ws = wb.active
        cName = self.getCurrentClass()
        row = 4
        col = column_index_from_string(config[cName]['excelCol'])
        while (ws.cell(row=row, column=col).value != None):
            row += 1
        #print(row)

        ws.cell(row=row, column=col).value = config[cName]['classComment']
        ws.cell(row=row, column=col).font = defaultFont


        wb.save(pathConfig.get('jindoPath',''))
        
    def inquiry(self):
        cur_date = QDate.currentDate()
        str_date = cur_date.toString("yyyy년 MM월 dd일 dddd")
        self.statusBar().showMessage(str_date)



app = QApplication(sys.argv)
window = MyWindow()
window.show()
app.exec_()